"""
src/output/export_ppt.py
==========================
Executive PPTX Report Generator

Produces a timestamped, board-ready PowerPoint from model outputs.
Uses python-pptx for PPTX generation.

Design rules:
  - Pure function: generate_pptx(report_data) → filepath
  - No financial calculations — receives pre-computed data
  - Timestamped filename: NVIDIA_FPA_Report_YYYYMMDD_HHMMSS.pptx
  - Output saved to outputs/ directory

Slides:
  1. Cover (company, scenario, date)
  2. Executive Summary (key KPIs + valuation headline)
  3. Financial Highlights (IS summary table)
  4. Valuation (DCF bridge + implied price)
  5. Scenario Comparison (Base vs Upside vs Downside delta)
  6. Appendix: Assumptions
"""

from __future__ import annotations

import logging
import os
from datetime import datetime

logger = logging.getLogger(__name__)

OUTPUT_DIR = "outputs"

# NVIDIA brand colours
_NAVY = "1B3A6B"
_GREEN = "76B900"  # NVIDIA green
_WHITE = "FFFFFF"
_GRAY = "F5F5F5"
_DARK = "333333"


def generate_pptx(
    report_data: dict,
    output_dir: str = OUTPUT_DIR,
    filename: str | None = None,
) -> str:
    """
    Generate a board-ready PPTX report from model outputs.

    Args:
        report_data: dict containing:
            - income_statement: pd.DataFrame
            - dcf_valuation:    dict
            - kpis:             dict (optional)
            - scenario:         str
            - run_timestamp:    str (optional)
        output_dir:  Directory to save the PPTX
        filename:    Optional custom filename (without extension)

    Returns:
        Absolute filepath of generated PPTX

    Raises:
        ImportError: If python-pptx is not installed
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX export. "
            "Install with: pip install python-pptx"
        )

    os.makedirs(output_dir, exist_ok=True)

    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    scenario = report_data.get("scenario", "base")
    if filename is None:
        filename = f"NVIDIA_FPA_Report_{scenario.upper()}_{timestamp}"
    filepath = os.path.join(output_dir, f"{filename}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    def _rgb(hex_str: str) -> RGBColor:
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _add_textbox(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=18,
        bold=False,
        color=_DARK,
        bg_color=None,
        align=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        if bg_color:
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(bg_color)
        return txBox

    def _add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(color)
        shape.line.fill.background()
        return shape

    def _add_table(
        slide, rows, left, top, width, height, header_bg=_NAVY, header_fg=_WHITE
    ):
        from pptx.util import Pt

        cols = len(rows[0])
        tbl = slide.shapes.add_table(
            len(rows), cols, Inches(left), Inches(top), Inches(width), Inches(height)
        ).table
        col_width = Inches(width / cols)
        for c in range(cols):
            tbl.columns[c].width = col_width

        for r, row_data in enumerate(rows):
            for c, cell_text in enumerate(row_data):
                cell = tbl.cell(r, c)
                cell.text = str(cell_text)
                tf = cell.text_frame
                tf.paragraphs[0].font.size = Pt(10)
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(header_bg)
                    tf.paragraphs[0].font.color.rgb = _rgb(header_fg)
                    tf.paragraphs[0].font.bold = True
                else:
                    bg = _GRAY if r % 2 == 0 else _WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(bg)
                    tf.paragraphs[0].font.color.rgb = _rgb(_DARK)
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        return tbl

    # ── Slide 1: Cover ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.333, 7.5, _NAVY)
    _add_textbox(
        slide,
        0.8,
        1.5,
        11,
        1.2,
        "NVIDIA Corporation",
        font_size=40,
        bold=True,
        color=_WHITE,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        0.8,
        2.8,
        11,
        0.8,
        "FP&A Performance Dashboard",
        font_size=28,
        bold=False,
        color=_GREEN,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        0.8,
        3.8,
        11,
        0.5,
        f"Scenario: {scenario.upper()}  |  Generated: {timestamp}  |  USD millions",
        font_size=14,
        color=_WHITE,
        align=PP_ALIGN.CENTER,
    )

    # ── Slide 2: Executive Summary ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.333, 0.8, _NAVY)
    _add_textbox(
        slide,
        0.3,
        0.1,
        10,
        0.6,
        "Executive Summary",
        font_size=22,
        bold=True,
        color=_WHITE,
    )

    dcf = report_data.get("dcf_valuation", {})
    kpis = report_data.get("kpis", {})

    # DCF headline
    price = dcf.get("implied_share_price_usd", 0)
    mkt = dcf.get("market_price_usd", 0)
    updown = dcf.get("upside_downside_pct", 0)
    ev = dcf.get("enterprise_value_usdm", 0)
    wacc_u = dcf.get("wacc_used", 0)
    terminal_growth_rate = dcf.get("terminal_growth_rate", 0)

    _add_textbox(
        slide,
        0.5,
        1.0,
        5.5,
        0.5,
        "Intrinsic Value Analysis",
        font_size=14,
        bold=True,
        color=_NAVY,
    )
    summary_rows = [
        ["Metric", "Value"],
        ["Implied Share Price", f"${price:.2f}"],
        ["Market Price", f"${mkt:.2f}"],
        ["Upside / (Downside)", f"{updown*100:.1f}%"],
        ["Enterprise Value", f"${ev:,.0f}M"],
        ["WACC", f"{wacc_u*100:.2f}%"],
        ["Terminal Growth", f"{terminal_growth_rate*100:.1f}%"],
    ]
    _add_table(slide, summary_rows, 0.5, 1.6, 5.5, 3.5)

    # KPI box
    if kpis:
        _add_textbox(
            slide,
            7.0,
            1.0,
            5.5,
            0.5,
            "Key Operating KPIs",
            font_size=14,
            bold=True,
            color=_NAVY,
        )
        kpi_display = [
            ["KPI", "Value"],
            [
                "Gross Margin",
                (
                    f"{kpis.get('gross_margin', 0)*100:.1f}%"
                    if kpis.get("gross_margin")
                    else "—"
                ),
            ],
            [
                "EBITDA Margin",
                (
                    f"{kpis.get('ebitda_margin', 0)*100:.1f}%"
                    if kpis.get("ebitda_margin")
                    else "—"
                ),
            ],
            [
                "FCF Margin",
                (
                    f"{kpis.get('fcf_margin', 0)*100:.1f}%"
                    if kpis.get("fcf_margin")
                    else "—"
                ),
            ],
            [
                "AR Days",
                f"{kpis.get('ar_days', 0):.1f} days" if kpis.get("ar_days") else "—",
            ],
            [
                "Current Ratio",
                (
                    f"{kpis.get('current_ratio', 0):.2f}x"
                    if kpis.get("current_ratio")
                    else "—"
                ),
            ],
        ]
        _add_table(slide, kpi_display, 7.0, 1.6, 5.5, 3.0)

    # ── Slide 3: Financial Highlights ────────────────────────────────────
    is_df = report_data.get("income_statement")
    if is_df is not None and not is_df.empty:
        slide = prs.slides.add_slide(blank_layout)
        _add_rect(slide, 0, 0, 13.333, 0.8, _NAVY)
        _add_textbox(
            slide,
            0.3,
            0.1,
            10,
            0.6,
            "Income Statement — 5-Year Projection",
            font_size=22,
            bold=True,
            color=_WHITE,
        )

        def _fmt(val, col):
            if val is None:
                return "—"
            if "pct" in str(col) or "margin" in str(col):
                return f"{float(val)*100:.1f}%"
            return f"{float(val):,.0f}"

        show_cols = [
            "fiscal_year",
            "revenue_usdm",
            "gross_profit_usdm",
            "ebitda_usdm",
            "net_income_usdm",
            "gross_margin_pct",
        ]
        show_labels = ["FY", "Revenue", "Gross Profit", "EBITDA", "Net Income", "GM%"]
        available = [
            (l, c) for l, c in zip(show_labels, show_cols) if c in is_df.columns
        ]

        tbl_rows = [[lbl for lbl, _ in available]]
        for _, row in is_df.iterrows():
            tbl_rows.append([_fmt(row.get(col), col) for _, col in available])

        _add_table(slide, tbl_rows, 0.5, 1.0, 12.3, 4.0)

    # ── Slide 4: Assumptions ─────────────────────────────────────────────
    assumptions = report_data.get("assumptions_used", {})
    if assumptions:
        slide = prs.slides.add_slide(blank_layout)
        _add_rect(slide, 0, 0, 13.333, 0.8, _NAVY)
        _add_textbox(
            slide,
            0.3,
            0.1,
            10,
            0.6,
            "Key Model Assumptions",
            font_size=22,
            bold=True,
            color=_WHITE,
        )

        wacc_a = assumptions.get("wacc", {})
        dcf_a = assumptions.get("dcf", {})
        assump_rows = [
            ["Assumption", "Value"],
            ["WACC", f"{wacc_a.get('wacc', 0)*100:.4f}%"],
            ["Risk-free Rate", f"{wacc_a.get('risk_free_rate', 0)*100:.2f}%"],
            ["Equity Risk Premium", f"{wacc_a.get('equity_risk_premium', 0)*100:.1f}%"],
            ["Beta (Blume adj.)", f"{wacc_a.get('beta_blume_adjusted', 0):.4f}"],
            ["Terminal Growth", f"{dcf_a.get('terminal_growth_rate', 0)*100:.1f}%"],
            ["Forecast Horizon", f"{dcf_a.get('projection_horizon_years', 5)} years"],
            [
                "Diluted Shares",
                f"{dcf_a.get('diluted_shares_outstanding_millions', 0):,.0f}M",
            ],
            ["Scenario", scenario.upper()],
        ]
        _add_table(slide, assump_rows, 0.5, 1.0, 6.0, 5.0)

    prs.save(filepath)
    logger.info("[export_pptx] PPTX generated: %s", filepath)
    return filepath
